"""Phase 3 — course content & ingestion tests (see docs/features/content-ingestion.md).

Covers: upload returns immediately as `processing` and enqueues; the pipeline produces
chunks with embeddings and flips to `indexed`; per-slide citation locations; a corrupt file
becomes `failed` and is surfaced; retrieval is course-scoped (NFR-5); bad type rejected;
non-owner blocked; delete removes file + chunks; reindex re-queues a failed document.

These run against the docker-compose Postgres + Redis (like the auth/courses tests).
"""

import io
import uuid

import pytest
from fastapi.testclient import TestClient
from sqlalchemy import select

from app.config import settings
from app.db import SessionLocal
from app.ingestion.embeddings import get_embedder
from app.ingestion.pipeline import process_document
from app.ingestion.queue import dequeue
from app.main import app
from app.models.content import Chunk
from app.redis_client import redis_client
from tests.conftest import TEST_DOMAIN


@pytest.fixture(autouse=True)
def clear_ingest_queue():
    """Isolate the queue so a stray job (e.g. from a live worker) can't leak into a test."""
    redis_client.delete(settings.ingest_queue_key)
    yield
    redis_client.delete(settings.ingest_queue_key)


# --- helpers ---------------------------------------------------------------

def _email() -> str:
    return f"user-{uuid.uuid4().hex[:10]}@{TEST_DOMAIN}"


def _login(sender, email: str) -> TestClient:
    client = TestClient(app, follow_redirects=False)
    client.post("/auth/request", json={"email": email})
    client.get(f"/auth/verify?token={sender.last_token()}")
    return client


def _instructor(monkeypatch, sender) -> TestClient:
    email = _email()
    monkeypatch.setattr(settings, "instructor_emails", email)
    return _login(sender, email)


def _make_course(client: TestClient, code="CS 310", name="OS") -> dict:
    resp = client.post("/courses", json={"code": code, "name": name})
    assert resp.status_code == 201, resp.text
    return resp.json()


def _upload(client, course_id, filename, data, content_type="application/octet-stream"):
    return client.post(
        f"/courses/{course_id}/documents",
        files={"file": (filename, data, content_type)},
    )


def _process_next() -> None:
    """Pop the next queued document id and run the pipeline (stands in for the worker)."""
    doc_id = dequeue(timeout=2)
    assert doc_id is not None, "expected a queued document"
    db = SessionLocal()
    try:
        process_document(db, doc_id)
    finally:
        db.close()


def _long_text(n=200) -> bytes:
    return ("Operating systems schedule processes and manage memory. " * n).encode("utf-8")


def _pptx_bytes(slide_texts) -> bytes:
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[5]  # title only
    for text in slide_texts:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- tests -----------------------------------------------------------------

def test_upload_returns_processing_and_enqueues(monkeypatch, sender):
    instr = _instructor(monkeypatch, sender)
    course = _make_course(instr)

    resp = _upload(instr, course["id"], "notes.txt", b"hello world", "text/plain")
    assert resp.status_code == 201, resp.text
    body = resp.json()
    assert body["status"] == "processing"
    assert body["chunk_count"] == 0
    assert body["filename"] == "notes.txt"
    assert body["type"] == "txt"

    # The id was queued for the worker; pop it to confirm and keep the queue clean.
    queued = dequeue(timeout=2)
    assert str(queued) == body["id"]


def test_pipeline_indexes_with_chunks(monkeypatch, sender):
    instr = _instructor(monkeypatch, sender)
    course = _make_course(instr)

    resp = _upload(instr, course["id"], "lecture.txt", _long_text(), "text/plain")
    doc_id = resp.json()["id"]
    _process_next()

    docs = instr.get(f"/courses/{course['id']}/documents").json()
    doc = next(d for d in docs if d["id"] == doc_id)
    assert doc["status"] == "indexed"
    assert doc["chunk_count"] > 1  # long text splits into several overlapping chunks
    assert doc["error"] is None

    # Chunks carry the course_id (NFR-5 scoping) and have embeddings.
    db = SessionLocal()
    try:
        chunks = db.execute(
            select(Chunk).where(Chunk.document_id == uuid.UUID(doc_id))
        ).scalars().all()
        assert len(chunks) == doc["chunk_count"]
        assert all(c.course_id == uuid.UUID(course["id"]) for c in chunks)
        assert all(len(c.embedding) == settings.embedding_dim for c in chunks)
    finally:
        db.close()


def test_pptx_locations_become_citations(monkeypatch, sender):
    instr = _instructor(monkeypatch, sender)
    course = _make_course(instr)

    data = _pptx_bytes(["Deadlock has four conditions", "Mutual exclusion and hold and wait"])
    resp = _upload(instr, course["id"], "week6.pptx", data)
    doc_id = resp.json()["id"]
    _process_next()

    db = SessionLocal()
    try:
        locations = {
            c.location
            for c in db.execute(
                select(Chunk).where(Chunk.document_id == uuid.UUID(doc_id))
            ).scalars().all()
        }
    finally:
        db.close()
    assert "slide 1" in locations
    assert "slide 2" in locations


def test_corrupt_file_marked_failed(monkeypatch, sender):
    instr = _instructor(monkeypatch, sender)
    course = _make_course(instr)

    resp = _upload(instr, course["id"], "broken.pdf", b"%PDF-1.4 not really a pdf", "application/pdf")
    doc_id = resp.json()["id"]
    _process_next()

    docs = instr.get(f"/courses/{course['id']}/documents").json()
    doc = next(d for d in docs if d["id"] == doc_id)
    assert doc["status"] == "failed"
    assert doc["error"]  # a human-readable reason is surfaced


def test_retrieval_is_course_scoped(monkeypatch, sender):
    instr = _instructor(monkeypatch, sender)
    c1 = _make_course(instr, code="CS 310", name="OS")
    c2 = _make_course(instr, code="CS 240", name="DS")

    _upload(instr, c1["id"], "os.txt", b"paging virtual memory and page tables " * 50, "text/plain")
    _process_next()
    _upload(instr, c2["id"], "ds.txt", b"binary trees and hash tables and graphs " * 50, "text/plain")
    _process_next()

    query_vec = get_embedder().embed_one("how does paging work")
    db = SessionLocal()
    try:
        scoped = db.execute(
            select(Chunk)
            .where(Chunk.course_id == uuid.UUID(c1["id"]))
            .order_by(Chunk.embedding.cosine_distance(query_vec))
            .limit(5)
        ).scalars().all()
        c2_chunks = db.execute(
            select(Chunk).where(Chunk.course_id == uuid.UUID(c2["id"]))
        ).scalars().all()
    finally:
        db.close()

    assert scoped, "expected results from course 1"
    assert all(c.course_id == uuid.UUID(c1["id"]) for c in scoped)
    assert c2_chunks, "course 2 should have its own chunks (so exclusion is meaningful)"


def test_bad_extension_rejected(monkeypatch, sender):
    instr = _instructor(monkeypatch, sender)
    course = _make_course(instr)
    resp = _upload(instr, course["id"], "malware.exe", b"MZ...", "application/octet-stream")
    assert resp.status_code == 422


def test_non_owner_cannot_upload_or_list(monkeypatch, sender):
    instr = _instructor(monkeypatch, sender)
    course = _make_course(instr)

    student = _login(sender, _email())
    assert _upload(student, course["id"], "x.txt", b"hi", "text/plain").status_code == 403
    assert student.get(f"/courses/{course['id']}/documents").status_code == 403


def test_delete_removes_document_and_chunks(monkeypatch, sender):
    instr = _instructor(monkeypatch, sender)
    course = _make_course(instr)
    resp = _upload(instr, course["id"], "del.txt", _long_text(), "text/plain")
    doc_id = resp.json()["id"]
    _process_next()

    assert instr.delete(f"/documents/{doc_id}").status_code == 204
    assert instr.get(f"/courses/{course['id']}/documents").json() == []

    db = SessionLocal()
    try:
        remaining = db.execute(
            select(Chunk).where(Chunk.document_id == uuid.UUID(doc_id))
        ).scalars().all()
        assert remaining == []
    finally:
        db.close()


def test_reindex_requeues_failed_document(monkeypatch, sender):
    instr = _instructor(monkeypatch, sender)
    course = _make_course(instr)
    resp = _upload(instr, course["id"], "broken.pdf", b"%PDF junk", "application/pdf")
    doc_id = resp.json()["id"]
    _process_next()  # → failed

    again = instr.post(f"/documents/{doc_id}/reindex")
    assert again.status_code == 200
    assert again.json()["status"] == "processing"
    assert again.json()["error"] is None

    queued = dequeue(timeout=2)
    assert str(queued) == doc_id  # it was re-enqueued
