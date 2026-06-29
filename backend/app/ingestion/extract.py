"""Step 1 of the pipeline: turn a raw file's bytes into labelled text segments.

Each segment carries a human-readable `location` (e.g. "p.12" or "slide 4") which is
preserved through chunking and ends up as the answer's citation. Returning segments per
page/slide is what makes precise citations possible.
"""

from __future__ import annotations

import io
from dataclasses import dataclass


@dataclass
class Segment:
    location: str  # e.g. "p.3", "slide 5"
    text: str


class ExtractionError(Exception):
    """Raised when a file can't be parsed (corrupt, wrong type, empty)."""


def extract_segments(filename: str, data: bytes, file_type: str) -> list[Segment]:
    """Dispatch on file type and return non-empty text segments."""
    ext = file_type.lower().lstrip(".")
    try:
        if ext == "pdf":
            segments = _extract_pdf(data)
        elif ext == "docx":
            segments = _extract_docx(data)
        elif ext == "pptx":
            segments = _extract_pptx(data)
        elif ext in ("txt", "md"):
            segments = _extract_text(data)
        else:
            raise ExtractionError(f"Unsupported file type: {ext}")
    except ExtractionError:
        raise
    except Exception as exc:  # noqa: BLE001 — any parser failure becomes a clear error
        raise ExtractionError(f"Could not read {filename}: {exc}") from exc

    cleaned = [Segment(s.location, s.text.strip()) for s in segments if s.text.strip()]
    if not cleaned:
        raise ExtractionError("No extractable text found")
    return cleaned


def _extract_pdf(data: bytes) -> list[Segment]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return [Segment(f"p.{i}", page.extract_text() or "") for i, page in enumerate(reader.pages, 1)]


def _extract_docx(data: bytes) -> list[Segment]:
    from docx import Document as DocxDocument

    doc = DocxDocument(io.BytesIO(data))
    # Word has no fixed pages at the XML level; treat the whole body as one segment.
    text = "\n".join(p.text for p in doc.paragraphs)
    return [Segment("document", text)]


def _extract_pptx(data: bytes) -> list[Segment]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    segments: list[Segment] = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [shape.text for shape in slide.shapes if shape.has_text_frame]
        segments.append(Segment(f"slide {i}", "\n".join(parts)))
    return segments


def _extract_text(data: bytes) -> list[Segment]:
    return [Segment("document", data.decode("utf-8", errors="replace"))]
